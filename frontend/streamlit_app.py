# frontend/streamlit_app.py
import streamlit as st
import pandas as pd
import sqlite3
import os
from io import BytesIO
import requests
from datetime import datetime
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from sqlalchemy import create_engine, text
from PIL import Image
import tempfile

st.set_page_config(page_title="Agent AI Ecom Analytics", layout="wide")

# -----------------------------
# Sidebar / config
# -----------------------------
DEFAULT_BACKEND = os.getenv("BACKEND_URL", "http://localhost:5000")

BACKEND_URL = st.sidebar.text_input("Backend URL", value=DEFAULT_BACKEND)

st.sidebar.header("Data source")
data_source = st.sidebar.selectbox("Select data source", options=["CSV (upload)", "Database (Postgres/MySQL)"])
db_conn_string = None
if data_source.startswith("Database"):
    db_conn_string = st.sidebar.text_input("SQLAlchemy connection string", placeholder="postgresql://user:pass@host:port/dbname")
    st.sidebar.caption("DB credentials are used only from this frontend instance (no server). Ensure network access.")

st.sidebar.header("Charts & Filters")
show_category_charts = st.sidebar.checkbox("Always show category charts from uploaded DB/CSV", value=True)
enable_category_filter = st.sidebar.checkbox("Enable category filter for queries", value=True)

# -----------------------------
# Helpers
# -----------------------------
def _post_generate_sql(nl_query, schema):
    try:
        resp = requests.post(f"{BACKEND_URL.rstrip('/')}/generate_sql", json={"nl_query": nl_query, "schema": schema}, timeout=60)
    except Exception as e:
        return None, f"Backend request failed: {e}"
    if resp.status_code != 200:
        try:
            return None, resp.json().get("error", resp.text)
        except Exception:
            return None, resp.text
    return resp.json(), None

def to_excel_bytes(df: pd.DataFrame) -> bytes:
    output = BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        df.to_excel(writer, index=False, sheet_name="results")
    return output.getvalue()

def df_to_sqlite_in_memory(df: pd.DataFrame, table_name="orders"):
    con = sqlite3.connect(":memory:")
    df.columns = [c.strip().replace(" ", "_") for c in df.columns]
    df.to_sql(table_name, con, index=False, if_exists="replace")
    return con

def export_pptx(result_df: pd.DataFrame, figures: list, title="Ecom Analytics Report"):
    prs = Presentation()
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
    except Exception:
        pass
    slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
    s = prs.slides.add_slide(slide_layout)
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    rows = min(30, len(result_df) + 1)
    cols = len(result_df.columns) if len(result_df.columns) > 0 else 1
    try:
        table_shape = s.shapes.add_table(rows, cols, left, top, width, Inches(3)).table
        for j, col in enumerate(result_df.columns):
            table_shape.cell(0, j).text = str(col)
        for i, (_, row) in enumerate(result_df.head(29).iterrows(), start=1):
            for j, col in enumerate(result_df.columns):
                table_shape.cell(i, j).text = str(row[col])
    except Exception:
        pass

    for fig in figures:
        try:
            img_temp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            fig.savefig(img_temp.name, bbox_inches='tight')
            img_temp.flush()
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_temp.name, Inches(1), Inches(1), width=Inches(8))
            img_temp.close()
        except Exception:
            continue

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()

def export_pdf(result_df: pd.DataFrame, figures: list, title="Ecom Analytics Report"):
    bio = BytesIO()
    c = canvas.Canvas(bio, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 16)
    c.drawString(40, height - 50, title)
    c.setFont("Helvetica", 10)
    c.drawString(40, height - 70, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    y = height - 100
    cols = list(result_df.columns)
    max_rows = min(20, len(result_df))
    col_width = (width - 80) / max(1, len(cols)) if cols else width - 80
    c.setFont("Helvetica-Bold", 9)
    for i, col in enumerate(cols):
        c.drawString(40 + i * col_width, y, str(col)[:15])
    y -= 12
    c.setFont("Helvetica", 8)
    for _, row in result_df.head(max_rows).iterrows():
        for i, col in enumerate(cols):
            text = str(row[col])[:15]
            c.drawString(40 + i * col_width, y, text)
        y -= 12
        if y < 100:
            c.showPage()
            y = height - 100
    for fig in figures:
        try:
            img_temp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            fig.savefig(img_temp.name, bbox_inches='tight')
            img_temp.flush()
            c.showPage()
            img = Image.open(img_temp.name)
            img_w, img_h = img.size
            factor = (width - 80) / img_w
            img_w_scaled = img_w * factor
            img_h_scaled = img_h * factor
            c.drawImage(ImageReader(img_temp.name), 40, height - img_h_scaled - 50, width=img_w_scaled, height=img_h_scaled)
            img_temp.close()
        except Exception:
            continue
    c.save()
    bio.seek(0)
    return bio.getvalue()

# -----------------------------
# Load data (CSV or DB)
# -----------------------------
df = None
engine = None
base_con = None

if data_source == "CSV (upload)":
    uploaded = st.file_uploader("Upload orders/transactions CSV", type=["csv"])
    if not uploaded:
        st.info("Upload a CSV (OrderID, CustomerName, ProductID, ProductName, ProductCategory, Quantity, Price, OrderDate recommended).")
        st.stop()
    try:
        df = pd.read_csv(uploaded)
    except Exception as e:
        try:
            uploaded.seek(0)
            df = pd.read_csv(uploaded, encoding="latin1")
        except Exception as e2:
            st.error(f"Failed to read uploaded CSV: {e} / {e2}")
            st.stop()
else:
    # Database selected
    if not db_conn_string:
        st.info("Enter a SQLAlchemy connection string in the sidebar to use a DB source.")
        st.stop()
    try:
        engine = create_engine(db_conn_string, pool_pre_ping=True)
        preview_q = "SELECT * FROM orders LIMIT 50"
        df = pd.read_sql_query(preview_q, engine)
    except Exception as exc:
        st.error(f"Error connecting to DB or fetching preview: {exc}")
        st.stop()

if df is None:
    st.error("Failed to load any data. Please check your CSV or DB connection.")
    st.stop()

st.success(f"Loaded dataset — {df.shape[0]} rows, {df.shape[1]} columns")
st.dataframe(df.head(50))

# detect date column
def detect_date_col(df_local):
    for c in df_local.columns:
        if 'date' in c.lower() or 'order' in c.lower():
            try:
                pd.to_datetime(df_local[c])
                return c
            except Exception:
                continue
    return None

date_col = detect_date_col(df)

# prepare sqlite in-memory if CSV
if data_source == "CSV (upload)":
    try:
        base_con = df_to_sqlite_in_memory(df.copy(), "orders")
    except Exception as e:
        st.error(f"Failed to create in-memory DB: {e}")
        st.stop()

# categories
categories = []
if 'ProductCategory' in df.columns:
    categories = sorted(df['ProductCategory'].dropna().unique().tolist())

if enable_category_filter and categories:
    selected_categories = st.sidebar.multiselect("Filter by ProductCategory", options=categories, default=categories)
else:
    selected_categories = categories

# -----------------------------
# UI: Chat + Results
# -----------------------------
if "history" not in st.session_state:
    st.session_state.history = []

left, right = st.columns([1, 2])
with left:
    st.header("Ask the AI (natural language)")
    nl_query = st.text_input("Example: 'Sankar how many products purchased this month?'")
    run_button = st.button("Generate SQL & Run")
    st.markdown("### Query history")
    for h in reversed(st.session_state.history[-20:]):
        st.write(f"**Q:** {h['nl']}")
        st.code(h['sql'])

with right:
    st.header("Result")
    if run_button:
        if not nl_query or nl_query.strip() == "":
            st.error("Type a natural language query.")
        else:
            st.info("Generating SQL from backend...")
            schema = {col: str(df[col].dtype) for col in df.columns}
            resp_json, err = _post_generate_sql(nl_query, schema)
            if err:
                st.error(f"Backend error: {err}")
                st.stop()
            sql_text = resp_json.get("sql")
            if not sql_text:
                st.error("Backend did not return SQL.")
                st.stop()
            st.markdown("**Generated SQL**")
            st.code(sql_text, language="sql")

            # Execute SQL
            result_df = None
            try:
                if data_source == "CSV (upload)":
                    con_to_use = base_con
                    if enable_category_filter and selected_categories and 'ProductCategory' in df.columns:
                        filtered_df = df[df['ProductCategory'].isin(selected_categories)].copy()
                        con_to_use = df_to_sqlite_in_memory(filtered_df, "orders")
                    if con_to_use is None:
                        st.error("No in-memory DB available to run SQL against.")
                        st.stop()
                    result_df = pd.read_sql_query(sql_text, con_to_use)
                else:
                    if engine is None:
                        st.error("No DB engine available to run SQL against.")
                        st.stop()
                    result_df = pd.read_sql_query(text(sql_text.rstrip(";")), engine)
            except Exception as exc:
                st.error(f"SQL execution error: {exc}")
                st.stop()

            if result_df is None:
                st.error("Query returned no DataFrame.")
                st.stop()

            st.session_state.history.append({"nl": nl_query, "sql": sql_text, "ts": datetime.now().isoformat()})
            st.markdown("### Query result (first 200 rows)")
            st.dataframe(result_df.head(200))

            # Build figures safely
            figures = []

            # Product distribution
            prod_col_candidates = [c for c in result_df.columns if 'product' in c.lower() and 'name' not in c.lower()]
            prod_col = prod_col_candidates[0] if prod_col_candidates else None
            if not prod_col and 'ProductName' in result_df.columns:
                prod_col = 'ProductName'
            if prod_col and result_df.shape[0] > 0:
                try:
                    pie_data = result_df[prod_col].value_counts().rename_axis(prod_col).reset_index(name='count')
                    fig, ax = plt.subplots()
                    pie_data_top = pie_data.head(10)
                    ax.pie(pie_data_top['count'], labels=pie_data_top[prod_col], autopct='%1.1f%%', startangle=90)
                    ax.axis('equal')
                    st.pyplot(fig)
                    figures.append(fig)
                except Exception:
                    pass

            # Category charts
            category_col = None
            if 'ProductCategory' in result_df.columns:
                category_col = 'ProductCategory'
            elif 'ProductID' in result_df.columns and 'ProductCategory' in df.columns:
                merged = result_df.merge(df[['ProductID', 'ProductCategory']].drop_duplicates(), on='ProductID', how='left')
                if 'ProductCategory' in merged.columns:
                    result_df['ProductCategory'] = merged['ProductCategory']
                    category_col = 'ProductCategory'

            if (category_col is None) and show_category_charts and 'ProductCategory' in df.columns:
                base_df = df.copy()
                if enable_category_filter and selected_categories:
                    base_df = base_df[base_df['ProductCategory'].isin(selected_categories)]
                try:
                    base_df['Revenue'] = pd.to_numeric(base_df['Price'], errors='coerce') * pd.to_numeric(base_df['Quantity'], errors='coerce')
                    cat_sales = base_df.groupby('ProductCategory')['Revenue'].sum().sort_values(ascending=False)
                    st.subheader("Sales by Category (uploaded dataset)")
                    st.bar_chart(cat_sales)
                    fig2 = cat_sales.plot(kind='bar', legend=False).get_figure()
                    figures.append(fig2)
                except Exception:
                    pass
            else:
                if category_col:
                    try:
                        result_df['Price'] = pd.to_numeric(result_df.get('Price', 0), errors='coerce').fillna(0)
                        result_df['Quantity'] = pd.to_numeric(result_df.get('Quantity', 1), errors='coerce').fillna(0)
                        result_df['Revenue'] = result_df['Price'] * result_df['Quantity']
                        cat_sales = result_df.groupby(category_col)['Revenue'].sum().sort_values(ascending=False)
                        st.subheader("Sales by Category (from query result)")
                        st.bar_chart(cat_sales)
                        fig2 = cat_sales.plot(kind='bar', legend=False).get_figure()
                        figures.append(fig2)
                    except Exception:
                        pass

            # Time series chart (best-effort)
            try:
                timeseries_col = None
                for c in result_df.columns:
                    if np.issubdtype(result_df[c].dtype, np.datetime64) or 'date' in c.lower() or 'time' in c.lower():
                        try:
                            result_df[c] = pd.to_datetime(result_df[c], errors='coerce')
                            if result_df[c].notna().sum() > 0:
                                timeseries_col = c
                                break
                        except Exception:
                            continue
                if timeseries_col is None and date_col is not None:
                    join_keys = [k for k in ("OrderID", "order_id", "OrderId", "Orderid") if k in result_df.columns and k in df.columns]
                    if join_keys:
                        key = join_keys[0]
                        merged = result_df.merge(df[[key, date_col]], on=key, how='left')
                        merged[date_col] = pd.to_datetime(merged[date_col], errors='coerce')
                        if merged[date_col].notna().sum() > 0:
                            result_df[date_col] = merged[date_col]
                            timeseries_col = date_col
                if timeseries_col:
                    st.subheader("Purchases over time")
                    period = st.selectbox("Aggregation period", options=["D","W","M"], index=2, key="agg_period")
                    metric = st.selectbox("Metric", options=["count_rows","sum_quantity","sum_revenue"], index=0, key="agg_metric")
                    tmp = result_df.copy()
                    tmp[timeseries_col] = pd.to_datetime(tmp[timeseries_col], errors='coerce')
                    tmp = tmp.dropna(subset=[timeseries_col]).set_index(timeseries_col)
                    qty_cols = [c for c in result_df.columns if c.lower() in ("quantity","qty","count","units")]
                    revenue_cols = [c for c in result_df.columns if 'price' in c.lower() or 'revenue' in c.lower()]
                    if metric == "count_rows":
                        ts = tmp.resample(period).size().rename("count")
                    elif metric == "sum_quantity" and qty_cols:
                        q = qty_cols[0]
                        tmp[q] = pd.to_numeric(tmp[q], errors='coerce').fillna(0)
                        ts = tmp.resample(period)[q].sum().rename("quantity_sum")
                    elif metric == "sum_revenue" and revenue_cols:
                        r = revenue_cols[0]
                        tmp[r] = pd.to_numeric(tmp[r], errors='coerce').fillna(0)
                        ts = tmp.resample(period)[r].sum().rename("revenue_sum")
                    else:
                        ts = tmp.resample(period).size().rename("count")
                    if not ts.empty:
                        st.line_chart(ts)
            except Exception:
                pass

            # Export buttons
            if result_df is not None and not result_df.empty:
                st.download_button("Download Excel", data=to_excel_bytes(result_df), file_name=f"query_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx")
                try:
                    pptx_bytes = export_pptx(result_df, figures, title=f"Report - {nl_query[:80]}")
                    st.download_button("Download PowerPoint (.pptx)", data=pptx_bytes, file_name=f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
                except Exception:
                    st.warning("PowerPoint export failed for this result.")
                try:
                    pdf_bytes = export_pdf(result_df, figures, title=f"Report - {nl_query[:80]}")
                    st.download_button("Download PDF", data=pdf_bytes, file_name=f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf")
                except Exception:
                    st.warning("PDF export failed for this result.")
